"""
Generate a non-technical stakeholder presentation for the HR Onboarding
Coordination buildable specification.

Design language (Claude-inspired):
- Minimal, typography-first
- Single warm accent (#D97757 — Claude's signature coral)
- Deep ink (#1F1F1E) on warm off-white (#F5F4EE)
- Large titles, short body, generous whitespace
- No clip-art, no gradients, no shadows

Inputs:  ../spec-hr-onboarding-coordination.md (content source, human-curated here)
Output:  ./HR-Onboarding-Coordination-Stakeholder-Overview.pptx

Run:  python build_stakeholder_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# -----------------------------------------------------------------------------
# Design tokens
# -----------------------------------------------------------------------------
INK        = RGBColor(0x1F, 0x1F, 0x1E)   # primary text
MUTED      = RGBColor(0x6B, 0x6A, 0x66)   # secondary text
ACCENT     = RGBColor(0xD9, 0x77, 0x57)   # Claude coral
CANVAS     = RGBColor(0xF5, 0xF4, 0xEE)   # warm off-white background
HAIRLINE   = RGBColor(0xE2, 0xDF, 0xD5)   # subtle rule

FONT_HEAD  = "Calibri"   # widely available; swap to "Inter" or "Söhne" if installed
FONT_BODY  = "Calibri"

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.8)
MARGIN_T = Inches(0.7)
CONTENT_W = SLIDE_W - Inches(1.6)

OUTPUT = Path(__file__).parent / "HR-Onboarding-Coordination-Stakeholder-Overview.pptx"


# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------
def set_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=INK, font=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=INK, font=FONT_BODY, line_spacing=1.3,
                bullet_char="—", bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(10)

        b = p.add_run()
        b.text = f"{bullet_char}   "
        b.font.name = font
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        b.font.bold = True

        r = p.add_run()
        r.text = item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_hairline(slide, left, top, width, color=HAIRLINE, weight=Pt(0.75)):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_accent_mark(slide, left, top, width=Inches(0.45), height=Inches(0.06)):
    mk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    mk.fill.solid()
    mk.fill.fore_color.rgb = ACCENT
    mk.line.fill.background()
    return mk


def add_slide_number(slide, n, total):
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
             Inches(0.9), Inches(0.3),
             f"{n} / {total}",
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text):
    add_text(slide, MARGIN_L, SLIDE_H - Inches(0.45),
             Inches(8), Inches(0.3),
             text, size=10, color=MUTED)


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------
def title_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CANVAS)
    add_accent_mark(s, MARGIN_L, Inches(2.2), width=Inches(0.9), height=Inches(0.08))
    add_text(s, MARGIN_L, Inches(2.5), CONTENT_W, Inches(1.4),
             "HR Onboarding Coordination",
             size=48, bold=True, color=INK, font=FONT_HEAD)
    add_text(s, MARGIN_L, Inches(3.5), CONTENT_W, Inches(0.8),
             "An agentic approach to a 220-hire-per-year workload",
             size=22, color=MUTED, font=FONT_HEAD)
    add_text(s, MARGIN_L, SLIDE_H - Inches(1.0), CONTENT_W, Inches(0.5),
             "Stakeholder overview  ·  First-draft proposal  ·  April 2026",
             size=12, color=MUTED)
    return s


def section_header(slide, kicker, title):
    add_text(slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(0.35),
             kicker.upper(), size=11, bold=True, color=ACCENT, font=FONT_HEAD)
    add_text(slide, MARGIN_L, MARGIN_T + Inches(0.35), CONTENT_W, Inches(0.9),
             title, size=32, bold=True, color=INK, font=FONT_HEAD)
    add_hairline(slide, MARGIN_L, MARGIN_T + Inches(1.25), Inches(1.5),
                 color=ACCENT, weight=Pt(1.5))


def problem_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "The problem", "Onboarding is a coordination job the team can't keep up with")

    # Three stat columns
    stats = [
        ("220+", "new hires per year"),
        ("~40",  "tasks per onboarding, across 2 weeks"),
        ("6",    "separate systems per onboarding"),
    ]
    col_w = (CONTENT_W - Inches(0.8)) / 3
    top = Inches(2.5)
    for i, (big, label) in enumerate(stats):
        x = MARGIN_L + i * (col_w + Inches(0.4))
        add_text(s, x, top, col_w, Inches(1.2), big,
                 size=56, bold=True, color=ACCENT, font=FONT_HEAD)
        add_text(s, x, top + Inches(1.3), col_w, Inches(0.8), label,
                 size=16, color=INK, font=FONT_BODY)

    add_text(s, MARGIN_L, Inches(5.2), CONTENT_W, Inches(1.5),
             "\u201CMost of this is paperwork my team should not be touching, "
             "but every time we try to automate, something falls through the cracks "
             "because the edge cases never look the same twice.\u201D",
             size=16, color=MUTED, font=FONT_HEAD)
    add_text(s, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.4),
             "— HR Ops lead, stated in scenario",
             size=12, color=MUTED)

    add_footer(s, "Source: Scenario 1, Week 1 Practice Scenario Pool")
    add_slide_number(s, idx, total)
    return s


def why_now_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "Why agentic, why now", "The shape of the work fits what agents do well — with limits")

    col_w = (CONTENT_W - Inches(0.6)) / 2
    top = Inches(2.3)

    # Left column — what fits
    add_text(s, MARGIN_L, top, col_w, Inches(0.4),
             "WHAT FITS", size=11, bold=True, color=ACCENT)
    add_bullets(s, MARGIN_L, top + Inches(0.5), col_w, Inches(4.5), [
        "~85% of tasks are routine coordination — moving data between known systems.",
        "~7,480 routine tasks per year. At that scale, a 3-person team cannot absorb the load.",
        "Tasks are rule-governed: IT tickets, benefits triggers, training assignments, scheduling.",
    ], size=15, line_spacing=1.35)

    # Right column — what doesn't
    x2 = MARGIN_L + col_w + Inches(0.6)
    add_text(s, x2, top, col_w, Inches(0.4),
             "WHAT DOESN'T", size=11, bold=True, color=ACCENT)
    add_bullets(s, x2, top + Inches(0.5), col_w, Inches(4.5), [
        "~15% of tasks are judgment calls — classification, buddy norms, I-9 hold decisions.",
        "These carry legal, compliance, and trust risk if an agent decides silently.",
        "They must stay human-led. The boundary between routine and judgment is the design.",
    ], size=15, line_spacing=1.35)

    add_footer(s, "85 / 15 split as cited in the scenario; not yet validated with live workflow data.")
    add_slide_number(s, idx, total)
    return s


def proposal_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "The proposal", "Three capabilities, one delegation boundary")

    top = Inches(2.3)
    caps = [
        ("1", "Onboarding Orchestrator",
         "Creates the onboarding record per hire, spawns the ~40 tasks, drives each "
         "to completion or to a human when judgment is required."),
        ("2", "Buddy Match Proposer",
         "Proposes up to 5 candidate buddies with ranked rationale. HR Ops confirms. "
         "The agent never assigns on its own."),
        ("3", "Compliance Training Assigner",
         "Enrols new hires in LMS courses based on role and jurisdiction. Tracks "
         "completion. Escalates when a role spans unclear jurisdictions."),
    ]
    row_h = Inches(1.4)
    for i, (num, title, body) in enumerate(caps):
        y = top + i * (row_h + Inches(0.15))
        # number block
        add_text(s, MARGIN_L, y, Inches(0.8), row_h,
                 num, size=44, bold=True, color=ACCENT, font=FONT_HEAD,
                 anchor=MSO_ANCHOR.TOP)
        # title + body
        add_text(s, MARGIN_L + Inches(0.9), y, CONTENT_W - Inches(0.9), Inches(0.4),
                 title, size=18, bold=True, color=INK, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(0.9), y + Inches(0.45),
                 CONTENT_W - Inches(0.9), Inches(0.9),
                 body, size=13, color=MUTED, font=FONT_BODY, line_spacing=1.35)

    add_slide_number(s, idx, total); add_footer(s, "Full spec: spec-hr-onboarding-coordination.md")
    return s


def boundary_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "The delegation boundary",
                   "The agent is a coordinator, not a decision-maker")

    col_w = (CONTENT_W - Inches(0.6)) / 2
    top = Inches(2.3)

    add_text(s, MARGIN_L, top, col_w, Inches(0.4),
             "THE AGENT HANDLES", size=11, bold=True, color=ACCENT)
    add_bullets(s, MARGIN_L, top + Inches(0.5), col_w, Inches(4.5), [
        "IT provisioning tickets for standard role-based bundles.",
        "Benefits-enrolment triggers once classification is set.",
        "Standard compliance-training enrolment and completion tracking.",
        "Welcome pack, calendar coordination, manager handoff summary at day 10.",
        "Detecting late I-9s, overdue tasks, and day-14 incompletes.",
    ], size=14, line_spacing=1.35)

    x2 = MARGIN_L + col_w + Inches(0.6)
    add_text(s, x2, top, col_w, Inches(0.4),
             "HUMANS ALWAYS DECIDE", size=11, bold=True, color=ACCENT)
    add_bullets(s, x2, top + Inches(0.5), col_w, Inches(4.5), [
        "Contractor vs. full-employee classification (legal exposure).",
        "Whether a late I-9 triggers a work-authorisation hold (IRCA).",
        "Whether a buddy pairing crosses seniority norms (social judgment).",
        "Non-standard IT requests (executive kit, accommodations, dev environments).",
        "Signing off that an onboarding is complete.",
    ], size=14, line_spacing=1.35)

    add_footer(s, "Every human decision is logged with user, reason, and timestamp — 7-year retention.")
    add_slide_number(s, idx, total)
    return s


def metrics_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "What success looks like",
                   "Four measurable outcomes — three targets, one hard rule")

    top = Inches(2.3)
    rows = [
        ("≥ 68%", "of all onboarding tasks executed without HR Ops touching them",
         "Target — needs validation"),
        ("≥ 50%", "reduction in HR Ops time per onboarding",
         "Target — baseline needed"),
        ("≤ 1%",  "of onboardings with any incomplete task at day 14",
         "Target — baseline needed"),
        ("100%",  "of judgment cases reach a named human before any decision is recorded",
         "Non-negotiable — this is the boundary"),
    ]
    row_h = Inches(1.05)
    for i, (big, label, qualifier) in enumerate(rows):
        y = top + i * row_h
        add_text(s, MARGIN_L, y, Inches(2.2), row_h, big,
                 size=34, bold=True, color=ACCENT, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(2.4), y + Inches(0.1), Inches(7.5), Inches(0.5),
                 label, size=16, color=INK, font=FONT_BODY)
        add_text(s, MARGIN_L + Inches(2.4), y + Inches(0.55), Inches(7.5), Inches(0.4),
                 qualifier, size=11, color=MUTED, font=FONT_BODY)
        if i < len(rows) - 1:
            add_hairline(s, MARGIN_L, y + row_h - Inches(0.05),
                         CONTENT_W, color=HAIRLINE)

    add_slide_number(s, idx, total)
    return s


def assumptions_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "What we still need to validate",
                   "Three questions we need answered before we build past the first pass")

    top = Inches(2.4)
    qs = [
        ("01", "The six systems",
         "The scenario names four: Workday, ServiceNow, the LMS, and email. "
         "We need the other two named before benefits tasks can be built."),
        ("02", "Who sets employment classification",
         "The agent cannot set contractor vs. employee. We need to confirm "
         "HR Ops owns that decision in Workday — not Legal, not Payroll."),
        ("03", "Is the 80% routine-delegation target realistic",
         "Walking the 40-task template with HR Ops would tell us how often "
         "'routine' tasks actually need a human eye in practice."),
    ]
    row_h = Inches(1.35)
    for i, (num, head, body) in enumerate(qs):
        y = top + i * row_h
        add_text(s, MARGIN_L, y, Inches(0.9), row_h,
                 num, size=26, bold=True, color=ACCENT, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(1.0), y, CONTENT_W - Inches(1.0), Inches(0.4),
                 head, size=17, bold=True, color=INK, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(1.0), y + Inches(0.45),
                 CONTENT_W - Inches(1.0), Inches(0.85),
                 body, size=13, color=MUTED, font=FONT_BODY, line_spacing=1.35)

    add_footer(s, "Full assumption log (9 items with hypotheses and confidence) in the spec.")
    add_slide_number(s, idx, total)
    return s


def risks_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "What could go wrong",
                   "And how the design absorbs it")

    rows = [
        ("An integration goes down",
         "The agent queues work, retries on a schedule, and flags Ops if the outage persists. Nothing is silently dropped."),
        ("A judgment case is misrouted as routine",
         "Every task carries a delegation class. Judgment classes cannot be marked complete by the agent — only by a logged human action."),
        ("A human doesn't respond in time",
         "Every escalation has a named owner and an SLA. Overdue items surface in a day-14 audit to the HR Ops lead."),
        ("The spec is wrong somewhere",
         "A closed build loop against an AI coding agent is run this week to pressure-test the spec. Gaps surface as diagnosable mismatches, not production bugs."),
    ]
    top = Inches(2.3)
    row_h = Inches(1.0)
    for i, (head, body) in enumerate(rows):
        y = top + i * row_h
        add_text(s, MARGIN_L, y + Inches(0.1), Inches(0.35), Inches(0.35),
                 "→", size=18, bold=True, color=ACCENT)
        add_text(s, MARGIN_L + Inches(0.45), y, Inches(4.0), Inches(0.4),
                 head, size=15, bold=True, color=INK, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(4.6), y, CONTENT_W - Inches(4.6), Inches(0.9),
                 body, size=13, color=MUTED, font=FONT_BODY, line_spacing=1.35)
        if i < len(rows) - 1:
            add_hairline(s, MARGIN_L, y + row_h - Inches(0.1),
                         CONTENT_W, color=HAIRLINE)

    add_slide_number(s, idx, total)
    return s


def next_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    section_header(s, "What we're asking for", "Three things, in order")

    steps = [
        ("Confirm the two unnamed systems",
         "A 15-minute conversation with the HR Ops lead unlocks the benefits-enrolment capability."),
        ("Validate one success-metric baseline",
         "We need at least one current-state number (time per onboarding OR incomplete-at-day-14 rate) to anchor the business case."),
        ("Green-light a 2-week build loop",
         "With the above, we can build and demo the Orchestrator capability end-to-end against a sample onboarding."),
    ]
    top = Inches(2.4)
    row_h = Inches(1.3)
    for i, (head, body) in enumerate(steps):
        y = top + i * row_h
        add_text(s, MARGIN_L, y, Inches(0.8), Inches(0.8),
                 f"{i+1}", size=46, bold=True, color=ACCENT, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(0.9), y, CONTENT_W - Inches(0.9), Inches(0.4),
                 head, size=18, bold=True, color=INK, font=FONT_HEAD)
        add_text(s, MARGIN_L + Inches(0.9), y + Inches(0.5),
                 CONTENT_W - Inches(0.9), Inches(0.7),
                 body, size=13, color=MUTED, font=FONT_BODY, line_spacing=1.35)

    add_slide_number(s, idx, total)
    return s


def closing_slide(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s, CANVAS)
    add_accent_mark(s, MARGIN_L, Inches(2.8), width=Inches(0.9), height=Inches(0.08))
    add_text(s, MARGIN_L, Inches(3.1), CONTENT_W, Inches(1.2),
             "Questions, reactions, objections.",
             size=40, bold=True, color=INK, font=FONT_HEAD)
    add_text(s, MARGIN_L, Inches(4.1), CONTENT_W, Inches(0.6),
             "The spec is a draft. Your pushback is what turns it into something buildable.",
             size=18, color=MUTED, font=FONT_HEAD)
    add_footer(s, "Full spec: Spec-Generation/spec-hr-onboarding-coordination.md")
    add_slide_number(s, idx, total)
    return s


# -----------------------------------------------------------------------------
# Build
# -----------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        title_slide,
        problem_slide,
        why_now_slide,
        proposal_slide,
        boundary_slide,
        metrics_slide,
        assumptions_slide,
        risks_slide,
        next_slide,
        closing_slide,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        build(prs, i, total)

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()

